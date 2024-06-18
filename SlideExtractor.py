from pptx.slide import Slide


class SlideExtractor:
    def __init__(self):
        pass
    def extract_text_from_slide(self, slide : Slide) -> str:
        text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text)
        return ' '.join(text).strip()
